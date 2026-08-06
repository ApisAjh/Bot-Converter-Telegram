"""
Converter Document → PDF: DOCX, XLSX, PPTX, TXT.

Sama seperti Word→PDF di converters/pdf.py, konversi di sini TIDAK memakai
LibreOffice — melainkan mengekstrak konten (teks/tabel/slide) lalu
merender ulang ke PDF dengan reportlab. Cocok untuk isi/dokumen berbasis
teks & tabel; layout visual asli (desain slide PPTX, styling XLSX, dsb)
tidak direplikasi persis. Lihat README bagian "Batasan".
"""
from __future__ import annotations

from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

from app.converters.pdf import word_to_pdf as _docx_to_pdf  # reuse existing logic

docx_to_pdf = _docx_to_pdf  # alias supaya konsisten dengan nama fitur


def xlsx_to_pdf(input_path: Path, output_path: Path) -> None:
    wb = load_workbook(str(input_path), data_only=True)
    styles = getSampleStyleSheet()
    story = []
    for idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        if idx > 0:
            story.append(PageBreak())
        story.append(Paragraph(f"Sheet: {sheet_name}", styles["Heading2"]))
        story.append(Spacer(1, 0.3 * cm))
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append(["" if v is None else str(v) for v in row])
        if data:
            table = Table(data)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
            ]))
            story.append(table)
    if not story:
        story.append(Paragraph("(Spreadsheet kosong)", styles["Normal"]))
    SimpleDocTemplate(str(output_path), pagesize=landscape(A4)).build(story)


def pptx_to_pdf(input_path: Path, output_path: Path) -> None:
    prs = Presentation(str(input_path))
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(prs.slides, start=1):
        if i > 1:
            story.append(PageBreak())
        story.append(Paragraph(f"Slide {i}", styles["Heading2"]))
        story.append(Spacer(1, 0.3 * cm))
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                story.append(Paragraph(shape.text_frame.text.strip(), styles["Normal"]))
                story.append(Spacer(1, 0.2 * cm))
    if not story:
        story.append(Paragraph("(Presentasi kosong)", styles["Normal"]))
    SimpleDocTemplate(str(output_path), pagesize=A4).build(story)


def txt_to_pdf(input_path: Path, output_path: Path) -> None:
    styles = getSampleStyleSheet()
    story = []
    text = input_path.read_text(encoding="utf-8", errors="replace")
    for line in text.splitlines():
        line = line.strip()
        if line:
            # escape karakter khusus reportlab/HTML-like markup
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, styles["Normal"]))
        else:
            story.append(Spacer(1, 0.2 * cm))
    if not story:
        story.append(Paragraph("(File teks kosong)", styles["Normal"]))
    SimpleDocTemplate(str(output_path), pagesize=A4).build(story)
